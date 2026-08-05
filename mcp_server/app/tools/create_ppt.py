"""Tool：產生 PowerPoint（.pptx）簡報。

用結構化參數（每張投影片的標題/條列/備忘稿/圖片），不是讓模型自己寫程式碼跑。
參數用 Pydantic model 宣告巢狀欄位，理由見 create_excel.py 開頭的說明。
"""
import base64
import io

from mcp_types import CallToolResult, TextContent
from pptx import Presentation
from pptx.util import Inches
from pydantic import BaseModel, Field


class SlideSpec(BaseModel):
    title: str = Field("", description="這張投影片的標題")
    bullets: list[str] = Field(default_factory=list, description="條列重點，依序列出")
    notes: str | None = Field(None, description="給講者看的備忘稿（選填）")
    image_base64: str | None = Field(None, description="這張投影片要放的圖片，base64 編碼內容（選填）")


def _build(slides: list[SlideSpec]) -> bytes:
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    for spec in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = spec.title
        if spec.bullets:
            body = slide.placeholders[1].text_frame
            body.text = spec.bullets[0]
            for b in spec.bullets[1:]:
                p = body.add_paragraph()
                p.text = b
        if spec.notes:
            slide.notes_slide.notes_text_frame.text = spec.notes
        if spec.image_base64:
            slide.shapes.add_picture(
                io.BytesIO(base64.b64decode(spec.image_base64)),
                Inches(5), Inches(1.5), width=Inches(4),
            )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def register(server) -> None:
    @server.tool(
        name="create_ppt",
        description="產生 PowerPoint（.pptx）簡報，每張投影片可有標題、條列重點、備忘稿、圖片。"
                    "何時使用：使用者要求做簡報、投影片、提案資料。",
    )
    def create_ppt(filename: str, slides: list[SlideSpec]) -> CallToolResult:
        name = filename if filename.endswith(".pptx") else filename + ".pptx"
        data = _build(slides)
        return CallToolResult(
            content=[TextContent(type="text", text=f"已產生簡報：{name}")],
            structured_content={
                "filename": name,
                "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "data_base64": base64.b64encode(data).decode(),
            },
        )
